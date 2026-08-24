from __future__ import annotations

import html
import json
import re
from datetime import datetime
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, Query
from fastapi.responses import FileResponse, HTMLResponse, PlainTextResponse
from sqlalchemy.orm import Session

from app.db.session import get_db
from app.models.core import AuditLog, ReportDraft

try:
    from docx import Document
    from docx.shared import Inches, Pt
except Exception:  # pragma: no cover
    Document = None  # type: ignore
    Inches = None  # type: ignore
    Pt = None  # type: ignore

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches as PptxInches, Pt as PptxPt
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore
    RGBColor = None  # type: ignore
    PptxInches = None  # type: ignore
    PptxPt = None  # type: ignore

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
except Exception:  # pragma: no cover
    A4 = None  # type: ignore
    getSampleStyleSheet = None  # type: ignore
    cm = None  # type: ignore
    Paragraph = None  # type: ignore
    SimpleDocTemplate = None  # type: ignore
    Spacer = None  # type: ignore

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
except Exception:  # pragma: no cover
    Workbook = None  # type: ignore
    Font = None  # type: ignore
    PatternFill = None  # type: ignore

router = APIRouter(prefix="/exports", tags=["exports"])

EXPORT_ROOT = Path(__file__).resolve().parents[1] / "storage" / "exports"
EXPORT_ROOT.mkdir(parents=True, exist_ok=True)

ALLOWED_FORMATS = {"docx", "pptx", "pdf", "xlsx", "txt", "html"}


def slugify(value: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "-", value.lower()).strip("-")
    return cleaned or "dalili-report"


def get_draft_or_404(db: Session, draft_id: int) -> ReportDraft:
    draft = db.query(ReportDraft).filter(ReportDraft.id == draft_id).first()
    if not draft:
        raise HTTPException(status_code=404, detail="Report draft not found")
    return draft


def get_latest_draft_or_404(db: Session, project_id: int | None, report_type: str | None) -> ReportDraft:
    query = db.query(ReportDraft)
    if project_id is not None:
        query = query.filter(ReportDraft.project_id == project_id)
    if report_type is not None:
        query = query.filter(ReportDraft.report_type == report_type)
    draft = query.order_by(ReportDraft.created_at.desc()).first()
    if not draft:
        raise HTTPException(status_code=404, detail="Report draft not found")
    return draft


def base_filename(draft: ReportDraft, extension: str) -> str:
    stamp = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
    return f"dalili-{draft.report_type}-{slugify(draft.title)}-{draft.id}-{stamp}.{extension}"


def report_lines(draft: ReportDraft) -> list[str]:
    lines = [line.rstrip() for line in (draft.report_text or "").splitlines()]
    return lines or [draft.title, "No report text was saved for this draft."]


def build_html(draft: ReportDraft) -> str:
    def paragraph(line: str) -> str:
        escaped = html.escape(line)
        if not line.strip():
            return "<br />"
        if line.isupper() and len(line) > 4:
            return f"<h1>{escaped}</h1>"
        if re.match(r"^\d+\.\s", line):
            return f"<h2>{escaped}</h2>"
        if line.startswith("- "):
            return f"<p class='bullet'>{escaped}</p>"
        return f"<p>{escaped}</p>"

    return f"""<!doctype html>
<html>
<head>
<meta charset="utf-8" />
<title>{html.escape(draft.title)}</title>
<style>
  body {{ font-family: Arial, Helvetica, sans-serif; color: #102033; line-height: 1.55; margin: 44px; }}
  .cover {{ background: #f2f4f7; border: 1px solid #e2e8f0; border-radius: 14px; padding: 18px; margin-bottom: 28px; }}
  .brand {{ color: #073B2A; font-size: 13px; font-weight: 800; letter-spacing: .16em; text-transform: uppercase; }}
  h1 {{ color: #073B2A; font-size: 24px; border-bottom: 3px solid #F5B400; padding-bottom: 10px; }}
  h2 {{ color: #073B2A; font-size: 18px; margin-top: 24px; }}
  p {{ font-size: 12.5px; margin: 7px 0; }}
  .bullet {{ margin-left: 18px; }}
  .footer {{ margin-top: 36px; color: #64748b; font-size: 11px; border-top: 1px solid #e2e8f0; padding-top: 12px; }}
</style>
</head>
<body>
<div class="cover">
  <div class="brand">Dalili report export</div>
  <h1>{html.escape(draft.title)}</h1>
  <p><strong>Report type:</strong> {html.escape(draft.report_type)}</p>
  <p><strong>Project ID:</strong> {draft.project_id}</p>
  <p><strong>Dataset:</strong> {html.escape(draft.file_name or 'Not specified')}</p>
  <p><strong>Indicator:</strong> {html.escape(draft.indicator_name or 'Not specified')}</p>
  <p><strong>Author:</strong> {html.escape(draft.author)}</p>
  <p><strong>Generated:</strong> {datetime.utcnow().isoformat(timespec='seconds')} UTC</p>
</div>
{''.join(paragraph(line) for line in report_lines(draft))}
<div class="footer">Generated by Dalili. Validate all calculations and source data before external submission.</div>
</body>
</html>"""


def make_txt(draft: ReportDraft) -> Path:
    path = EXPORT_ROOT / base_filename(draft, "txt")
    header = f"Dalili Report Export\nTitle: {draft.title}\nReport type: {draft.report_type}\nProject ID: {draft.project_id}\nGenerated: {datetime.utcnow().isoformat(timespec='seconds')} UTC\n\n"
    path.write_text(header + (draft.report_text or ""), encoding="utf-8")
    return path


def make_html(draft: ReportDraft) -> Path:
    path = EXPORT_ROOT / base_filename(draft, "html")
    path.write_text(build_html(draft), encoding="utf-8")
    return path


def make_docx(draft: ReportDraft) -> Path:
    if Document is None:
        raise HTTPException(status_code=500, detail="python-docx is not installed")
    path = EXPORT_ROOT / base_filename(draft, "docx")
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_heading(draft.title, level=0)
    title.runs[0].font.name = "Arial"
    title.runs[0].font.size = Pt(20)
    meta = doc.add_paragraph()
    meta.add_run("Dalili report export\n").bold = True
    meta.add_run(f"Report type: {draft.report_type}\n")
    meta.add_run(f"Project ID: {draft.project_id}\n")
    meta.add_run(f"Dataset: {draft.file_name or 'Not specified'}\n")
    meta.add_run(f"Indicator: {draft.indicator_name or 'Not specified'}\n")
    meta.add_run(f"Generated: {datetime.utcnow().isoformat(timespec='seconds')} UTC")

    for line in report_lines(draft):
        if not line.strip():
            doc.add_paragraph("")
        elif line.isupper() and len(line) > 4:
            doc.add_heading(line, level=1)
        elif re.match(r"^\d+\.\s", line):
            doc.add_heading(line, level=2)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.add_paragraph("Generated by Dalili. Validate all calculations and source data before external submission.")
    doc.save(path)
    return path


def make_pptx(draft: ReportDraft) -> Path:
    if Presentation is None:
        raise HTTPException(status_code=500, detail="python-pptx is not installed")
    path = EXPORT_ROOT / base_filename(draft, "pptx")
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_slide(title: str, bullets: list[str], big_number: str | None = None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(242, 244, 247)
        tx = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.4), PptxInches(12.1), PptxInches(0.8))
        tf = tx.text_frame
        tf.text = title
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = PptxPt(34)
        tf.paragraphs[0].font.color.rgb = RGBColor(7, 59, 42)
        if big_number:
            bx = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(1.55), PptxInches(4), PptxInches(1.2))
            btf = bx.text_frame
            btf.text = big_number
            btf.paragraphs[0].font.bold = True
            btf.paragraphs[0].font.size = PptxPt(54)
        body = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.7 if not big_number else 3.0), PptxInches(11.5), PptxInches(4.8))
        frame = body.text_frame
        frame.clear()
        for item in bullets[:7]:
            p = frame.add_paragraph()
            p.text = item[:220]
            p.level = 0
            p.font.size = PptxPt(22)
            p.font.color.rgb = RGBColor(16, 32, 51)

    lines = [line for line in report_lines(draft) if line.strip() and not set(line.strip()) <= {"="}]
    add_slide("Dalili results deck", [f"Report: {draft.title}", f"Dataset: {draft.file_name or 'Not specified'}", f"Indicator: {draft.indicator_name or 'Not specified'}", f"Generated: {datetime.utcnow().date().isoformat()}"])
    if draft.indicator_name:
        add_slide("Headline result", [f"Indicator: {draft.indicator_name}", f"Quality score: {draft.quality_score if draft.quality_score is not None else 'Not available'}", "Use the report draft for calculation details."], None)
    add_slide("Key findings", [line for line in lines if not re.match(r"^\d+\.\s", line)][:6])
    add_slide("Recommended actions", [line[2:] for line in lines if line.startswith("- ")][:7] or ["Validate source data.", "Review caveats and data quality issues.", "Confirm narrative before external sharing."])
    add_slide("Evidence and limitations", ["This deck was generated from a saved Dalili report draft.", "Validate all figures against the source dataset before donor submission.", f"Draft record ID: {draft.id}"])
    prs.save(path)
    return path


def make_pdf(draft: ReportDraft) -> Path:
    if SimpleDocTemplate is None:
        raise HTTPException(status_code=500, detail="reportlab is not installed")
    path = EXPORT_ROOT / base_filename(draft, "pdf")
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(f"<b>{html.escape(draft.title)}</b>", styles["Title"]))
    meta = f"Dalili report export<br/>Report type: {html.escape(draft.report_type)}<br/>Project ID: {draft.project_id}<br/>Dataset: {html.escape(draft.file_name or 'Not specified')}<br/>Indicator: {html.escape(draft.indicator_name or 'Not specified')}<br/>Generated: {datetime.utcnow().isoformat(timespec='seconds')} UTC"
    story.append(Paragraph(meta, styles["Normal"]))
    story.append(Spacer(1, 0.35 * cm))
    for line in report_lines(draft):
        if not line.strip():
            story.append(Spacer(1, 0.18 * cm))
        elif line.isupper() and len(line) > 4:
            story.append(Paragraph(f"<b>{html.escape(line)}</b>", styles["Heading1"]))
        elif re.match(r"^\d+\.\s", line):
            story.append(Paragraph(f"<b>{html.escape(line)}</b>", styles["Heading2"]))
        else:
            story.append(Paragraph(html.escape(line), styles["Normal"]))
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=1.5 * cm, leftMargin=1.5 * cm, topMargin=1.4 * cm, bottomMargin=1.4 * cm)
    doc.build(story)
    return path


def make_xlsx(draft: ReportDraft) -> Path:
    if Workbook is None:
        raise HTTPException(status_code=500, detail="openpyxl is not installed")
    path = EXPORT_ROOT / base_filename(draft, "xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    header_fill = PatternFill(start_color="073B2A", end_color="073B2A", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    rows = [
        ["Field", "Value"],
        ["Title", draft.title],
        ["Report type", draft.report_type],
        ["Project ID", draft.project_id],
        ["Dataset ID", draft.dataset_id or ""],
        ["Indicator result ID", draft.indicator_result_id or ""],
        ["File name", draft.file_name or ""],
        ["Indicator name", draft.indicator_name or ""],
        ["Quality score", draft.quality_score if draft.quality_score is not None else ""],
        ["Author", draft.author],
        ["Generated", datetime.utcnow().isoformat(timespec="seconds") + " UTC"],
    ]
    for row in rows:
        ws.append(row)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
    ws.column_dimensions["A"].width = 26
    ws.column_dimensions["B"].width = 90

    text_ws = wb.create_sheet("Report text")
    text_ws.append(["Line number", "Text"])
    for cell in text_ws[1]:
        cell.fill = header_fill
        cell.font = header_font
    for idx, line in enumerate(report_lines(draft), start=1):
        text_ws.append([idx, line])
    text_ws.column_dimensions["A"].width = 14
    text_ws.column_dimensions["B"].width = 120

    if draft.insight_summary_json:
        insight_ws = wb.create_sheet("Insight summary")
        insight_ws.append(["Key", "Value"])
        for cell in insight_ws[1]:
            cell.fill = header_fill
            cell.font = header_font
        try:
            summary = json.loads(draft.insight_summary_json)
            if isinstance(summary, dict):
                for key, value in summary.items():
                    insight_ws.append([str(key), json.dumps(value) if isinstance(value, (dict, list)) else str(value)])
        except Exception:
            insight_ws.append(["raw", draft.insight_summary_json])
        insight_ws.column_dimensions["A"].width = 28
        insight_ws.column_dimensions["B"].width = 90

    wb.save(path)
    return path


def make_export(draft: ReportDraft, export_format: str) -> tuple[Path, str, str]:
    fmt = export_format.lower().strip()
    if fmt not in ALLOWED_FORMATS:
        raise HTTPException(status_code=400, detail=f"Unsupported export format. Use one of: {', '.join(sorted(ALLOWED_FORMATS))}")
    if fmt == "docx":
        return make_docx(draft), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"
    if fmt == "pptx":
        return make_pptx(draft), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"
    if fmt == "pdf":
        return make_pdf(draft), "application/pdf", "pdf"
    if fmt == "xlsx":
        return make_xlsx(draft), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"
    if fmt == "html":
        return make_html(draft), "text/html; charset=utf-8", "html"
    return make_txt(draft), "text/plain; charset=utf-8", "txt"


def log_export(db: Session, draft: ReportDraft, export_format: str) -> None:
    db.add(AuditLog(action="export", entity_type="report_draft", entity_id=str(draft.id), detail=f"{draft.title} exported as {export_format}"))
    db.commit()


@router.get("/formats")
def list_export_formats():
    return {
        "formats": sorted(ALLOWED_FORMATS),
        "recommended": ["docx", "pptx", "pdf", "xlsx"],
        "note": "Exports are generated from saved backend report drafts.",
    }


@router.get("/report-drafts/latest")
def export_latest_report_draft(
    project_id: int | None = Query(default=None),
    report_type: str | None = Query(default=None),
    format: str = Query(default="docx"),
    db: Session = Depends(get_db),
):
    draft = get_latest_draft_or_404(db, project_id, report_type)
    path, media_type, _ = make_export(draft, format)
    log_export(db, draft, format)
    return FileResponse(path, media_type=media_type, filename=path.name)


@router.get("/report-drafts/{draft_id}")
def export_report_draft(
    draft_id: int,
    format: str = Query(default="docx"),
    db: Session = Depends(get_db),
):
    draft = get_draft_or_404(db, draft_id)
    path, media_type, _ = make_export(draft, format)
    log_export(db, draft, format)
    return FileResponse(path, media_type=media_type, filename=path.name)


@router.get("/report-drafts/{draft_id}/preview", response_class=HTMLResponse)
def preview_report_draft(draft_id: int, db: Session = Depends(get_db)):
    draft = get_draft_or_404(db, draft_id)
    return HTMLResponse(build_html(draft))


@router.get("/report-drafts/{draft_id}/text", response_class=PlainTextResponse)
def report_draft_text(draft_id: int, db: Session = Depends(get_db)):
    draft = get_draft_or_404(db, draft_id)
    return PlainTextResponse(draft.report_text or "")
